"""
sample.py — requirements.txt の各パッケージ使用例

実行方法:
    uv run python scripts/sample.py               # 全セクション実行
    uv run python scripts/sample.py path/to/file.pdf  # PDF を指定して変換

各セクションは独立した関数になっています。
使いたい部分だけ呼び出すことも可能です。
"""

import sys
from pathlib import Path

ROOT_DIR   = Path(__file__).parent.parent
OUTPUT_DIR = ROOT_DIR / "output"
DATA_DIR   = ROOT_DIR / "data"
OUTPUT_DIR.mkdir(exist_ok=True)
DATA_DIR.mkdir(exist_ok=True)


# ─────────────────────────────────────────────
# 1. NumPy — 数値計算
# ─────────────────────────────────────────────
def demo_numpy():
    import numpy as np

    print("\n── NumPy ──────────────────────────────────")

    arr = np.array([1, 2, 3, 4, 5])
    print(f"配列:        {arr}")
    print(f"合計:        {arr.sum()}")
    print(f"平均:        {arr.mean()}")
    print(f"標準偏差:    {arr.std():.4f}")

    # 2次元配列（行列）
    matrix = np.arange(1, 10).reshape(3, 3)
    print(f"\n3×3 行列:\n{matrix}")
    print(f"転置行列:\n{matrix.T}")


# ─────────────────────────────────────────────
# 2. pandas — データフレーム操作 / CSV 読み書き
# ─────────────────────────────────────────────
def demo_pandas():
    import pandas as pd

    print("\n── pandas ─────────────────────────────────")

    # DataFrame を辞書から作成
    data = {
        "名前":   ["Alice", "Bob", "Charlie", "Diana"],
        "年齢":   [25, 30, 22, 28],
        "都市":   ["東京", "大阪", "福岡", "札幌"],
        "スコア": [85.5, 92.0, 78.3, 95.1],
    }
    df = pd.DataFrame(data)
    print(df.to_string(index=False))

    # 基本統計
    print(f"\n平均スコア: {df['スコア'].mean():.1f}")
    print(f"最高スコア: {df['スコア'].max()}")

    # CSV に保存して読み直す
    csv_path = OUTPUT_DIR / "sample.csv"
    df.to_csv(csv_path, index=False, encoding="utf-8-sig")  # utf-8-sig で Excel でも文字化けしない
    df_loaded = pd.read_csv(csv_path, encoding="utf-8-sig")
    print(f"\nCSV 保存 → 読み込み成功: {csv_path}")
    print(f"行数: {len(df_loaded)}, 列数: {len(df_loaded.columns)}")


# ─────────────────────────────────────────────
# 3. requests — HTTP 通信
# ─────────────────────────────────────────────
def demo_requests():
    import requests

    print("\n── requests ───────────────────────────────")

    url = "https://httpbin.org/get"
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        data = response.json()
        print(f"ステータスコード: {response.status_code}")
        print(f"送信元 IP:       {data.get('origin', '不明')}")
        print(f"User-Agent:      {data['headers'].get('User-Agent', '不明')[:60]}")
    except requests.exceptions.ConnectionError:
        print("ネットワーク接続エラー（オフライン環境では正常）")
    except requests.exceptions.Timeout:
        print("タイムアウト")


# ─────────────────────────────────────────────
# 4. python-dotenv — 環境変数の読み込み
# ─────────────────────────────────────────────
def demo_dotenv():
    import os
    from dotenv import load_dotenv

    print("\n── python-dotenv ──────────────────────────")

    env_path = Path(__file__).parent.parent / ".env"
    if not env_path.exists():
        print(f".env ファイルが見つかりません: {env_path}")
        print("  → .env.sample をコピーして .env を作成してください。")
        print("     cp .env.sample .env")
        return

    load_dotenv(env_path)

    # 読み込んだ変数の存在確認（値は表示しない）
    keys = ["DATA_DIR", "SOME_API_KEY"]
    for key in keys:
        val = os.getenv(key)
        status = "✓ 設定済み" if val else "– 未設定"
        print(f"  {key}: {status}")


# ─────────────────────────────────────────────
# 5. openpyxl — Excel ファイルの読み書き
# ─────────────────────────────────────────────
def demo_openpyxl():
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    print("\n── openpyxl ───────────────────────────────")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "サンプル"

    # ヘッダー行（太字・背景色付き）
    headers = ["名前", "年齢", "都市", "スコア"]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(fill_type="solid", fgColor="4472C4")
        cell.alignment = Alignment(horizontal="center")

    # データ行
    rows = [
        ("Alice",   25, "東京", 85.5),
        ("Bob",     30, "大阪", 92.0),
        ("Charlie", 22, "福岡", 78.3),
        ("Diana",   28, "札幌", 95.1),
    ]
    for row_data in rows:
        ws.append(row_data)

    # 列幅を自動調整
    for col in ws.columns:
        max_len = max(len(str(cell.value or "")) for cell in col)
        ws.column_dimensions[col[0].column_letter].width = max_len + 4

    xlsx_path = OUTPUT_DIR / "sample.xlsx"
    wb.save(xlsx_path)
    print(f"Excel 作成: {xlsx_path}")

    # 読み直して確認
    wb2 = openpyxl.load_workbook(xlsx_path)
    ws2 = wb2.active
    print(f"シート名: {ws2.title}, 行数: {ws2.max_row - 1} 件（ヘッダー除く）")


# ─────────────────────────────────────────────
# 6. python-pptx — PowerPoint ファイルの作成
# ─────────────────────────────────────────────
def demo_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    print("\n── python-pptx ────────────────────────────")

    prs = Presentation()

    # スライド1：タイトルスライド
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "サンプルプレゼンテーション"
    slide1.placeholders[1].text = "python-pptx による自動生成"

    # スライド2：箇条書き
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "使用パッケージ"
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "データ処理"
    items = ["NumPy — 数値計算", "pandas — データフレーム", "openpyxl — Excel", "python-pptx — PowerPoint"]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    pptx_path = OUTPUT_DIR / "sample.pptx"
    prs.save(pptx_path)
    print(f"PowerPoint 作成: {pptx_path}")
    print(f"スライド数: {len(prs.slides)}")


# ─────────────────────────────────────────────
# 7. pymupdf4llm — PDF → Markdown 変換
#    デフォルトで内閣府 AI 指針 2025 をダウンロードして変換します。
#    引数でローカルの PDF パスを渡すことも可能です。
# ─────────────────────────────────────────────

AI_GUIDELINE_URL = (
    "https://www8.cao.go.jp/cstp/ai/ai_guideline/ai_gl_2025.pdf"
)
AI_GUIDELINE_PDF = DATA_DIR / "ai_gl_2025.pdf"


def download_pdf(url: str, dest: Path) -> bool:
    """URL から PDF をダウンロードして dest に保存する。成功したら True を返す。"""
    import requests

    print(f"ダウンロード中: {url}")
    try:
        resp = requests.get(url, timeout=60, stream=True)
        resp.raise_for_status()
        total = int(resp.headers.get("content-length", 0))
        downloaded = 0
        with open(dest, "wb") as f:
            for chunk in resp.iter_content(chunk_size=8192):
                f.write(chunk)
                downloaded += len(chunk)
                if total:
                    pct = downloaded / total * 100
                    print(f"\r  {pct:5.1f}%  ({downloaded:,} / {total:,} bytes)", end="", flush=True)
        print(f"\r  完了: {dest.name} ({downloaded:,} bytes)         ")
        return True
    except Exception as e:
        print(f"\nダウンロード失敗: {e}")
        return False


def demo_pymupdf4llm(pdf_path: str | None = None):
    import pymupdf4llm

    print("\n── pymupdf4llm ────────────────────────────")

    # PDF の決定：引数 > デフォルト（内閣府 AI 指針）
    if pdf_path:
        pdf_file = Path(pdf_path)
        if not pdf_file.exists():
            print(f"ファイルが見つかりません: {pdf_file}")
            return
    else:
        pdf_file = AI_GUIDELINE_PDF
        if not pdf_file.exists():
            ok = download_pdf(AI_GUIDELINE_URL, pdf_file)
            if not ok:
                return

    print(f"変換中: {pdf_file.name} ...")
    md_text = pymupdf4llm.to_markdown(str(pdf_file))

    md_path = OUTPUT_DIR / (pdf_file.stem + ".md")
    md_path.write_text(md_text, encoding="utf-8")

    print(f"出力: {md_path}")
    print(f"文字数: {len(md_text):,} 文字")
    print("\n── 先頭 500 文字 ──")
    print(md_text[:500])


# ─────────────────────────────────────────────
# エントリーポイント
# ─────────────────────────────────────────────
if __name__ == "__main__":
    pdf_arg = sys.argv[1] if len(sys.argv) > 1 else None

    demo_numpy()
    demo_pandas()
    demo_requests()
    demo_dotenv()
    demo_openpyxl()
    demo_pptx()
    demo_pymupdf4llm(pdf_arg)

    print(f"\n出力ファイルは {OUTPUT_DIR} に保存されました。")
